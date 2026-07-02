# -*- coding: utf-8 -*-
"""
Builds the Odisea project presentation (Bulgarian), Odisea-branded to match the
web app's design tokens. Reuses the banner, the rendered thesis diagrams and the
app screenshots. Minimal text per slide by design.

Run:  python docs/presentation/build_pptx.py
Out:  docs/presentation/Odisea_Prezentatsiya.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
ASSETS = os.path.join(ROOT, "docs", "assets")
DIAG = os.path.join(ROOT, "docs", "thesis", "diagrams")
SHOTS = os.path.join(ROOT, "docs", "thesis", "screenshots")
OUT = os.path.join(ROOT, "docs", "presentation", "Odisea_Prezentatsiya.pptx")

# ---- Brand palette (from frontend/portal/src/styles.css) --------------------
TEAL_900 = RGBColor(0x0B, 0x2A, 0x2E)
TEAL_800 = RGBColor(0x0F, 0x3A, 0x40)
TEAL_700 = RGBColor(0x14, 0x47, 0x4E)
TEAL_600 = RGBColor(0x1A, 0x5A, 0x61)
TEAL_50 = RGBColor(0xEE, 0xF6, 0xF6)
GOLD = RGBColor(0xC6, 0x9B, 0x54)
GOLD_600 = RGBColor(0xB3, 0x89, 0x3F)
GOLD_200 = RGBColor(0xEC, 0xDC, 0xBA)
GOLD_50 = RGBColor(0xFA, 0xF4, 0xE7)
INK = RGBColor(0x0E, 0x16, 0x18)
GRAY_700 = RGBColor(0x3A, 0x47, 0x44)
GRAY_600 = RGBColor(0x52, 0x5F, 0x5C)
GRAY_500 = RGBColor(0x6C, 0x79, 0x76)
GRAY_200 = RGBColor(0xE1, 0xE6, 0xE4)
GRAY_100 = RGBColor(0xF1, 0xF4, 0xF3)
BG = RGBColor(0xF7, 0xF9, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFB, 0xFC, 0xFB)

# Fonts: Georgia is the brand's declared serif fallback (Prata/Lora -> Georgia),
# guaranteed on Windows; Segoe UI stands in for the Onest UI sans.
HEAD = "Georgia"
BODY = "Segoe UI"

EMU_IN = 914400
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _noline(shape):
    shape.line.fill.background()


def rect(s, l, t, w, h, color, rounded=False, line=None, line_w=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    if line is None:
        _noline(shp)
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if rounded:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def bg(s, color):
    rect(s, 0, 0, SW, SH, color)


def text(s, l, t, w, h, runs, size=16, color=INK, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    for txt, ov in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(ov.get("size", size))
        r.font.name = ov.get("font", font)
        r.font.bold = ov.get("bold", bold)
        r.font.italic = ov.get("italic", italic)
        r.font.color.rgb = ov.get("color", color)
    return tb


def para(tb, runs, size=16, color=INK, font=BODY, bold=False, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, spacing=1.0, level=0):
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = spacing
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, ov in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(ov.get("size", size))
        r.font.name = ov.get("font", font)
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
    return p


def eyebrow(s, txt, l, t, color=GOLD_600):
    text(s, l, t, 8, 0.35, txt.upper(), size=12.5, color=color, bold=True, font=BODY)


def heading(s, txt, l, t, w, size=38, color=INK):
    text(s, l, t, w, 1.4, txt, size=size, color=color, font=HEAD, bold=True, spacing=1.02)


def gold_rule(s, l, t, w=1.1, weight=0.045):
    rect(s, l, t, w, weight, GOLD)


def bullets(s, items, l, t, w, size=16.5, gap=13, color=INK, marker=GOLD_600,
            lh=1.12):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = lh
        p.space_after = Pt(gap)
        r0 = p.add_run()
        r0.text = "•  "
        r0.font.size = Pt(size)
        r0.font.name = BODY
        r0.font.bold = True
        r0.font.color.rgb = marker
        r1 = p.add_run()
        r1.text = it
        r1.font.size = Pt(size)
        r1.font.name = BODY
        r1.font.color.rgb = color
    return tb


def footer(s, idx, dark=False):
    c = TEAL_50 if dark else GRAY_500
    text(s, 0.55, SH - 0.5, 4, 0.3,
         [("odisea", {"font": HEAD, "bold": True, "size": 12, "color": GOLD if dark else TEAL_600})],
         size=12)
    text(s, SW - 1.4, SH - 0.5, 0.85, 0.3, f"{idx:02d}", size=11, color=c,
         align=PP_ALIGN.RIGHT)


def fit(s, path, l, t, w, h, card=True, pad=0.16, border=GRAY_200):
    """Place an image contained within (l,t,w,h), centered, on a white card."""
    iw, ih = Image.open(path).size
    r = iw / ih
    bw, bh = w - 2 * pad, h - 2 * pad
    if bw / bh > r:
        dh, dw = bh, bh * r
    else:
        dw, dh = bw, bw / r
    dl = l + (w - dw) / 2
    dt = t + (h - dh) / 2
    if card:
        rect(s, l, t, w, h, WHITE, rounded=True, line=border, line_w=1.0)
    s.shapes.add_picture(path, Inches(dl), Inches(dt), Inches(dw), Inches(dh))


def card(s, l, t, w, h, fill=WHITE, line=GRAY_200, line_w=1.0, top_accent=None):
    rect(s, l, t, w, h, fill, rounded=True, line=line, line_w=line_w)
    if top_accent:
        rect(s, l + 0.0, t, w, 0.09, top_accent, rounded=True)


# =============================================================================
# SLIDES
# =============================================================================

# ---- 1. Cover ----
def s_cover():
    s = slide()
    bg(s, TEAL_900)
    rect(s, 0, 0, SW, 0.16, GOLD)
    # banner (has the gold leaf mark + "odisea" wordmark)
    banner = os.path.join(ASSETS, "odisea-banner.png")
    iw, ih = Image.open(banner).size
    r = iw / ih
    bw = 9.4
    bh = bw / r
    s.shapes.add_picture(banner, Inches((SW - bw) / 2), Inches(1.05), Inches(bw), Inches(bh))
    y = 1.05 + bh + 0.45
    text(s, 1, y, SW - 2, 0.6,
         "Платформа за агрегиране и вграждане на туристически оферти",
         size=21, color=WHITE, font=HEAD, align=PP_ALIGN.CENTER)
    text(s, 1, y + 0.72, SW - 2, 0.5,
         "Един каталог, много канали. Една оферта, много марки.",
         size=15, color=GOLD_200, font=BODY, align=PP_ALIGN.CENTER, italic=True)
    text(s, 1, SH - 0.72, SW - 2, 0.3,
         "Аркан Ахмедов · Магистърска разработка · 2026",
         size=12.5, color=TEAL_50, align=PP_ALIGN.CENTER)


# ---- generic content header ----
def head_block(s, tag, title, tw=8.6):
    eyebrow(s, tag, 0.7, 0.62)
    gold_rule(s, 0.72, 0.98)
    heading(s, title, 0.68, 1.12, tw, size=34)


# ---- 2. Problem ----
def s_problem(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Контекст", "Разединена екосистема")
    bullets(s, [
        "Операторите създават офертите; агенциите ги продават.",
        "Липсва общ технологичен слой между двете страни.",
        "Ръчно поддържаните списъци остаряват бързо.",
        "iframe вграждането чупи бранда и потребителския опит.",
        "Изборът досега: актуалност или собствена марка.",
    ], 0.72, 2.35, 6.0)
    fit(s, os.path.join(DIAG, "fig-4-1-conceptual.png"), 7.15, 1.95, 5.5, 4.6)
    footer(s, i)


# ---- 3. Solution ----
def s_solution(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Решение", "Публикувай веднъж, разпространявай навсякъде", tw=11.8)
    bullets(s, [
        "Централен каталог от туристически оферти.",
        "Агенцията подбира, стилизира и публикува.",
        "Вградими уеб компоненти вместо iframe.",
        "Компонентите приемат облика на приемащия сайт.",
        "Код за вграждане — copy/paste, без разработчик.",
    ], 0.72, 2.4, 6.1)
    fit(s, os.path.join(SHOTS, "landing.png"), 7.0, 2.15, 5.75, 4.4)
    footer(s, i)


# ---- 4. Why powerful ----
def s_power(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Стойност", "Синтез, а не компромис")
    cards = [
        ("Винаги свежо", "Централно управление на съдържанието."),
        ("Твоят бранд", "Локален облик чрез дизайн токени и теми."),
        ("Едни данни", "Много канали и визуални контексти."),
        ("Две публики", "No-code за агенции; API за разработчици (в план)."),
    ]
    x = 0.72
    cw, gap = 2.92, 0.18
    for idx, (t, d) in enumerate(cards):
        cx = x + idx * (cw + gap)
        card(s, cx, 2.5, cw, 2.7, fill=WHITE, top_accent=GOLD)
        rect(s, cx + 0.28, 2.85, 0.5, 0.5, TEAL_50, rounded=True)
        text(s, cx + 0.28, 2.9, 0.5, 0.4, str(idx + 1), size=17, color=TEAL_600,
             font=HEAD, bold=True, align=PP_ALIGN.CENTER)
        text(s, cx + 0.28, 3.6, cw - 0.5, 0.5, t, size=17.5, color=INK, font=HEAD, bold=True)
        text(s, cx + 0.28, 4.15, cw - 0.5, 1.0, d, size=13.5, color=GRAY_600, spacing=1.1)
    text(s, 0.72, 5.6, 11.8, 0.6,
         [("Не компромис между ръчно и iframe, ", {"color": GRAY_600, "size": 15}),
          ("а тяхна синтеза.", {"color": TEAL_600, "size": 15, "bold": True, "font": HEAD})],
         )
    footer(s, i)


# ---- 5. Competitor analysis ----
def s_competitors(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Пазар", "PeakView срещу Odisea")
    # left card — PeakView
    card(s, 0.72, 2.3, 5.7, 4.35, fill=WHITE, line=GRAY_200)
    text(s, 1.05, 2.6, 5, 0.4, "PeakView: доминиращ в България", size=16.5,
         color=GRAY_700, font=HEAD, bold=True)
    bullets(s, [
        "White label по оператор, държава, транспорт.",
        "XML интеграции и изработка на сайтове.",
        "Travel Web Office: документи, фактури, ваучери.",
        "iframe подход: слаба UX и липса на SEO.",
    ], 1.05, 3.2, 5.1, size=14.5, gap=10, color=GRAY_600)
    # right card — Odisea
    card(s, 6.9, 2.3, 5.7, 4.35, fill=TEAL_900, line=TEAL_900, top_accent=GOLD)
    text(s, 7.23, 2.6, 5, 0.4, "Odisea: нашето предимство", size=16.5,
         color=WHITE, font=HEAD, bold=True)
    bullets(s, [
        "Модерни уеб компоненти, не iframe.",
        "Дизайн токени, а не само цветове.",
        "Connector платформа, не отделни интеграции.",
        "Разделени Collection · Experience · Theme · Publication.",
        "No-code builder днес; developer API в пътната карта.",
    ], 7.23, 3.2, 5.1, size=14.5, gap=9, color=TEAL_50, marker=GOLD)
    footer(s, i)


# ---- 6. Key concepts (domain chain) ----
def s_concepts(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Модел", "От оферта до публикация")
    chain = [
        ("Offer", "нормализирана оферта", TEAL_600),
        ("Collection", "подбор с филтри", TEAL_600),
        ("Experience\n+ Theme", "изглед и стил", TEAL_700),
        ("Publication", "публичен ключ", GOLD_600),
        ("Канали", "сайт · WP · API", TEAL_800),
    ]
    n = len(chain)
    x0, y = 0.72, 2.7
    cw, gap = 2.15, 0.32
    for idx, (t, d, col) in enumerate(chain):
        cx = x0 + idx * (cw + gap)
        card(s, cx, y, cw, 1.75, fill=WHITE, line=GRAY_200, top_accent=col)
        text(s, cx + 0.12, y + 0.34, cw - 0.24, 0.7, t, size=16.5, color=INK,
             font=HEAD, bold=True, align=PP_ALIGN.CENTER, spacing=0.95)
        text(s, cx + 0.12, y + 1.16, cw - 0.24, 0.5, d, size=12, color=GRAY_600,
             align=PP_ALIGN.CENTER)
        if idx < n - 1:
            text(s, cx + cw - 0.02, y + 0.55, gap + 0.04, 0.6, "→", size=22,
                 color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    text(s, 0.72, 4.95, 11.9, 0.9,
         [("Едни и същи оферти протичат през подбор, композиция и стил, "
           "за да станат ", {"color": GRAY_600, "size": 15}),
          ("публикация с един стабилен публичен ключ.", {"color": TEAL_600, "size": 15, "bold": True})],
         spacing=1.25)
    # connectors + multi-tenant note strip
    for j, (lbl) in enumerate(["Connector engine · XML/JSON", "Мулти-тенант изолация", "Webhooks · Leads · Analytics"]):
        bx = 0.72 + j * 4.03
        rect(s, bx, 6.0, 3.85, 0.62, GRAY_100, rounded=True)
        text(s, bx + 0.25, 6.16, 3.5, 0.4, lbl, size=12.5, color=GRAY_700, bold=True)
    footer(s, i)


# ---- 7. Architecture ----
def s_architecture(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Архитектура", "Clean Architecture, три повърхности")
    layers = [("WebAPI", "контролери · middleware"), ("Infrastructure", "EF Core · конектори"),
              ("Application", "случаи на употреба · DTO"), ("Domain", "ядро без зависимости")]
    y = 2.4
    for idx, (t, d) in enumerate(layers):
        cy = y + idx * 0.82
        shade = [TEAL_600, TEAL_700, TEAL_800, TEAL_900][idx]
        rect(s, 0.72, cy, 5.2, 0.72, shade, rounded=True)
        text(s, 0.95, cy + 0.13, 3, 0.5, t, size=15.5, color=WHITE, font=HEAD, bold=True)
        text(s, 3.2, cy + 0.17, 2.6, 0.4, d, size=11.5, color=TEAL_50, align=PP_ALIGN.RIGHT)
    text(s, 0.72, 5.75, 5.2, 0.5, "Зависимости само навътре → Domain не знае за нищо.",
         size=12.5, color=GRAY_600, italic=True, spacing=1.1)
    fit(s, os.path.join(DIAG, "fig-4-3-c4-container.png"), 6.35, 2.25, 6.3, 4.35)
    footer(s, i)


# ---- 8. Database ----
def s_database(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Данни", "PostgreSQL: релации и jsonb")
    bullets(s, [
        "Мулти-тенант: organizations, memberships.",
        "offers и source_offers с проследимост.",
        "collections, themes, publications.",
        "supplier_connections, import_jobs.",
        "Enum-и като текст; jsonb за конфигурации.",
    ], 0.72, 2.4, 5.0, size=15.5)
    fit(s, os.path.join(DIAG, "fig-4-6-er.png"), 5.85, 1.95, 6.8, 4.7)
    footer(s, i)


# ---- 9. Platform flow ----
def s_flow(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Поток", "От данни до запитване")
    steps = [
        "Операторът внася оферти (ръчно / групово / конектор).",
        "Системата нормализира и съхранява.",
        "Агенцията подбира (Collection + филтри).",
        "Избира изглед и тема, после публикува.",
        "Компонентът зарежда манифеста по публичен ключ.",
        "Туристът праща запитване → Lead + webhook.",
    ]
    y = 2.4
    for idx, st in enumerate(steps):
        cy = y + idx * 0.66
        rect(s, 0.72, cy, 0.44, 0.44, TEAL_600, rounded=True)
        text(s, 0.72, cy + 0.05, 0.44, 0.36, str(idx + 1), size=15, color=WHITE,
             font=HEAD, bold=True, align=PP_ALIGN.CENTER)
        text(s, 1.32, cy + 0.06, 5.4, 0.5, st, size=14, color=INK, spacing=1.0)
    fit(s, os.path.join(DIAG, "fig-4-7-sequence-publish.png"), 7.0, 2.0, 5.7, 4.7)
    footer(s, i)


# ---- 10. Merchant onboarding ----
def s_onboarding(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Onboarding", "Три роли, три пътя")
    cols = [
        ("Туроператор", TEAL_600, [
            "Регистрация → профил оператор.",
            "Внос на оферти или XML конектор.",
            "Управление на връзки и история.",
        ]),
        ("Агенция без разработчик", GOLD_600, [
            "Избира вече внесени оферти.",
            "Collection builder с филтри.",
            "Готова тема → 1 клик „Публикувай“.",
            "Copy/paste код на сайта си.",
        ]),
        ("Агенция с разработчик", TEAL_800, [
            "API ключ и публични REST endpoints.",
            "Webhooks за lead събития.",
            "SDK и developer docs (в развитие).",
        ]),
    ]
    x = 0.72
    cw, gap = 3.9, 0.22
    for idx, (t, col, items) in enumerate(cols):
        cx = x + idx * (cw + gap)
        card(s, cx, 2.35, cw, 4.15, fill=WHITE, top_accent=col)
        text(s, cx + 0.3, 2.68, cw - 0.6, 0.6, t, size=16, color=INK, font=HEAD, bold=True)
        gold_rule(s, cx + 0.32, 3.28, 0.7)
        bullets(s, items, cx + 0.32, 3.5, cw - 0.62, size=13.5, gap=9, color=GRAY_700)
    footer(s, i)


# ---- 11. Embed security ----
def s_embed(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Вграждане", "С ключ: оферти. Без ключ: нищо.")
    bullets(s, [
        "Всяка агенция има публикуем ключ.",
        "Компонентът праща Authorization: ApiKey.",
        "Ключът е безопасен за публичен HTML.",
        "Заключване по разрешени домейни.",
    ], 0.72, 2.45, 5.5, size=15.5)
    badges = [("200", "Валиден ключ", TEAL_600, "Оферти се зареждат"),
              ("401", "Без ключ", GOLD_600, "Достъпът е отказан"),
              ("403", "Чужд ключ", RGBColor(0xB2, 0x1E, 0x3B), "Друга агенция")]
    y = 2.45
    for idx, (code, t, col, d) in enumerate(badges):
        cy = y + idx * 1.32
        card(s, 6.7, cy, 5.9, 1.15, fill=WHITE, line=GRAY_200)
        rect(s, 6.95, cy + 0.24, 1.15, 0.66, col, rounded=True)
        text(s, 6.95, cy + 0.36, 1.15, 0.45, code, size=19, color=WHITE, font=HEAD,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, 8.35, cy + 0.24, 4, 0.4, t, size=15.5, color=INK, font=HEAD, bold=True)
        text(s, 8.35, cy + 0.66, 4, 0.4, d, size=12.5, color=GRAY_600)
    footer(s, i)


# ---- 12. Achieved ----
def s_achieved(i):
    s = slide()
    bg(s, BG)
    head_block(s, "Статус", "Функционален MVP")
    items = [
        "Три повърхности: Landing, Portal, Components",
        "Clean Architecture backend",
        "PostgreSQL с миграции",
        "Мулти-тенант изолация",
        "Оферти: CRUD и групов внос",
        "Collection builder с филтри",
        "Theme система с дизайн токени",
        "Publication с публичен ключ",
        "Lit web components (grid, carousel)",
        "Connector engine (XML, JSON)",
        "Lead система с webhook известия",
        "Автоматизирани тестове + Docker",
    ]
    colw = 6.0
    for idx, it in enumerate(items):
        col = idx // 6
        row = idx % 6
        cx = 0.75 + col * 6.15
        cy = 2.4 + row * 0.72
        rect(s, cx, cy + 0.02, 0.34, 0.34, TEAL_50, rounded=True)
        text(s, cx, cy - 0.02, 0.34, 0.34, "✓", size=15, color=TEAL_600, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, cx + 0.5, cy + 0.0, colw - 0.6, 0.5, it, size=14, color=GRAY_700, spacing=1.0)
    footer(s, i)


# ---- 13. Product showcase ----
def s_showcase(i):
    s = slide()
    bg(s, TEAL_900)
    eyebrow(s, "Продукт", 0.7, 0.62, color=GOLD)
    gold_rule(s, 0.72, 0.98)
    heading(s, "Порталът в действие", 0.68, 1.12, 9, size=34, color=WHITE)
    trio = [("builder.png", "Композитор"), ("marketplace.png", "Пазар на теми"),
            ("connections.png", "Конектори")]
    x = 0.72
    cw, gap = 3.9, 0.22
    for idx, (img, cap) in enumerate(trio):
        cx = x + idx * (cw + gap)
        fit(s, os.path.join(SHOTS, img), cx, 2.35, cw, 3.1, card=True, border=TEAL_700)
        text(s, cx, 5.55, cw, 0.4, cap, size=14, color=TEAL_50, align=PP_ALIGN.CENTER,
             font=HEAD, bold=True)
    text(s, 0.72, 6.3, 11.9, 0.4,
         "Работи на живо в контейнеризирана среда (Docker · Render · PostgreSQL).",
         size=13, color=GOLD_200, align=PP_ALIGN.CENTER, italic=True)
    footer(s, i, dark=True)


# ---- 14. Closing ----
def s_closing():
    s = slide()
    bg(s, TEAL_900)
    rect(s, 0, 0, SW, 0.16, GOLD)
    text(s, 1, 2.5, SW - 2, 1.2, "odisea", size=72, color=WHITE, font=HEAD, bold=True,
         align=PP_ALIGN.CENTER)
    gold_rule(s, (SW - 1.4) / 2, 3.95, 1.4)
    text(s, 1, 4.2, SW - 2, 0.6,
         "Един каталог, много канали. Една оферта, много марки.",
         size=19, color=GOLD_200, font=HEAD, align=PP_ALIGN.CENTER, italic=True)
    text(s, 1, 5.1, SW - 2, 0.5,
         "Избери съдържание, избери преживяване, избери стил и публикувай навсякъде.",
         size=14, color=TEAL_50, align=PP_ALIGN.CENTER)
    text(s, 1, SH - 0.9, SW - 2, 0.4, "Благодаря!", size=16, color=WHITE,
         align=PP_ALIGN.CENTER, bold=True)


# ---- build ----
s_cover()
s_problem(2)
s_solution(3)
s_power(4)
s_competitors(5)
s_concepts(6)
s_architecture(7)
s_database(8)
s_flow(9)
s_onboarding(10)
s_embed(11)
s_achieved(12)
s_showcase(13)
s_closing()

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
